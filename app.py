import streamlit as st
from rake_nltk import Rake
import nltk
from nltk.corpus import stopwords
import os
import io
from PyPDF2 import PdfReader
import docx
from pptx import Presentation
import re
import base64

nltk.download('stopwords')
stop_words = set(stopwords.words('english'))

try:
    nltk.data.find('tokenizers/punkt')
except LookupError:
    nltk.download('punkt')
    
def extract_keywords(text):
    r = Rake()
    r.extract_keywords_from_text(text)
    return r.get_ranked_phrases()

def clean_and_filter_keywords(keywords):
    unique_keywords = set()
    filtered_keywords = []
    for keyword in keywords:
        keyword_cleaned = re.sub(r'[^a-zA-Z0-9\s]', '', keyword).lower()
        if len(keyword_cleaned) > 4 and keyword_cleaned not in stop_words:
            if len(keyword_cleaned.split()) <= 2:
                if keyword_cleaned not in unique_keywords:
                    unique_keywords.add(keyword_cleaned)
                    filtered_keywords.append(keyword_cleaned)
    return filtered_keywords

def extract_text_from_pdf(pdf_bytes):
    pdf_file = io.BytesIO(pdf_bytes)
    pdf_reader = PdfReader(pdf_file)
    text = ""
    for page in pdf_reader.pages:
        text += page.extract_text()
    return text

def extract_text_from_docx(docx_bytes):
    docx_file = io.BytesIO(docx_bytes)
    doc = docx.Document(docx_file)
    text = ""
    for paragraph in doc.paragraphs:
        text += paragraph.text + " "
    return text

def extract_text_from_pptx(pptx_bytes):
    prs = Presentation(io.BytesIO(pptx_bytes))
    text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + " "
    return text

def main():
    st.title("PDF Extraction Web App")
    st.sidebar.header("Upload a PDF File")

    uploaded_file = st.sidebar.file_uploader("(PDF, DOCX, PPTX, or TXT)", type=["pdf", "docx", "pptx", "txt"])

    if uploaded_file is not None:
        st.sidebar.write("File uploaded successfully.")
        
        file_bytes = uploaded_file.read()

        file_extension = uploaded_file.name.split(".")[-1]

        if file_extension == "pdf":
            text = extract_text_from_pdf(file_bytes)
        elif file_extension == "docx":
            text = extract_text_from_docx(file_bytes)
        elif file_extension == "pptx":
            text = extract_text_from_pptx(file_bytes)
        elif file_extension == "txt":
            text = file_bytes.decode('utf-8') 
        else:
            st.write("Unsupported file type. Please upload a PDF, DOCX, PPTX, or TXT file.")
            return

        keywords = extract_keywords(text)

        filtered_keywords = clean_and_filter_keywords(keywords)

        st.subheader("File Contents:")
        st.write(text)

        st.subheader("Extracted Keywords:")
        for keyword in filtered_keywords:
            st.write(keyword)

        if st.button("Download Keywords File"):
            # Prepare the content of the keywords file
            keywords_content = "\n".join(filtered_keywords)

            # Set up the file name for download
            file_name = "extracted_keywords.txt"

            # Create a download link for the keywords file
            href = f"data:text/plain;charset=utf-8,{keywords_content.encode('utf-8').decode('latin-1')}"
            st.markdown(f'<a href="{href}" download="{file_name}">Click here to download the keywords file</a>', unsafe_allow_html=True)

if __name__ == "__main__":
    main()
